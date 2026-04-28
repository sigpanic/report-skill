import os
import tempfile
import logging
from typing import Optional, cast

logger = logging.getLogger(__name__)


def parse_course_material(file_path: str, enable_ocr: bool = True) -> dict:
    ext = os.path.splitext(file_path)[1].lower()

    if ext == '.pptx':
        return _parse_pptx(file_path, enable_ocr)
    elif ext == '.ppt':
        return _parse_ppt(file_path, enable_ocr)
    elif ext == '.docx':
        return _parse_docx(file_path, enable_ocr)
    elif ext == '.doc':
        return _parse_doc(file_path, enable_ocr)
    else:
        return {"error": f"不支持的文件格式: {ext}", "text": ""}


def _parse_pptx(file_path: str, enable_ocr: bool = True) -> dict:
    try:
        from pptx import Presentation
        prs = Presentation(file_path)

        slides = []
        ocr_texts = []
        for i, slide in enumerate(prs.slides):
            slide_content = {
                "slide_number": i + 1,
                "texts": [],
                "has_images": False
            }
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_content["texts"].append(text)
                if shape.shape_type == 13:
                    slide_content["has_images"] = True
                    if enable_ocr:
                        ocr_text = _ocr_image_from_shape(shape, i + 1)
                        if ocr_text:
                            ocr_texts.append(f"[幻灯片{i+1}图片文字]: {ocr_text}")
                if shape.has_table:
                    from pptx.table import Table as PptxTable
                    table = cast(PptxTable, shape.table)
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            row_text.append(cell.text.strip())
                        slide_content["texts"].append(" | ".join(row_text))
            slides.append(slide_content)

        all_text = []
        has_images = False
        for slide in slides:
            all_text.extend(slide["texts"])
            if slide.get("has_images"):
                has_images = True

        if ocr_texts:
            all_text.extend(ocr_texts)

        image_note = None
        if has_images:
            if ocr_texts:
                image_note = f"课件中包含图片，已通过OCR提取部分文字（{len(ocr_texts)}处）。如OCR结果不准确，请用户提供图片中的文字内容。"
            else:
                image_note = "课件中包含图片，图片内容无法自动提取。如果图片中包含重要文字信息，请用户手动提供图片中的文字内容。"

        return {
            "file_path": file_path,
            "format": "pptx",
            "total_slides": len(slides),
            "has_images": has_images,
            "ocr_enabled": enable_ocr and has_images,
            "slides": slides,
            "full_text": "\n".join(all_text),
            "image_note": image_note
        }
    except Exception as e:
        return {"error": str(e), "text": ""}


def _ocr_image_from_shape(shape, slide_num: int) -> Optional[str]:
    try:
        import pytesseract
        from PIL import Image
        import io

        image = shape.image
        image_bytes = image.blob
        img = Image.open(io.BytesIO(image_bytes))

        text = pytesseract.image_to_string(img, lang='chi_sim+eng')
        return text.strip() if text.strip() else None
    except ImportError:
        logger.info("pytesseract not installed, OCR skipped")
        return None
    except Exception as e:
        logger.debug("OCR failed for slide %d: %s", slide_num, e)
        return None


def _parse_ppt(file_path: str, enable_ocr: bool = True) -> dict:
    powerpoint = None
    presentation = None
    try:
        import win32com.client
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Visible = False
        powerpoint.DisplayAlerts = False

        abs_path = os.path.abspath(file_path)
        presentation = powerpoint.Presentations.Open(abs_path)

        tmp_dir = tempfile.mkdtemp()
        pptx_path = os.path.join(tmp_dir, "converted.pptx")
        presentation.Saveas(pptx_path, 24)
        presentation.Close()
        presentation = None
        powerpoint.Quit()
        powerpoint = None

        result = _parse_pptx(pptx_path, enable_ocr)
        result["format"] = "ppt (converted)"
        return result
    except Exception as e:
        if presentation:
            try:
                presentation.Close()
            except Exception:
                pass
        if powerpoint:
            try:
                powerpoint.Quit()
            except Exception:
                pass
        try:
            result = _parse_pptx(file_path, enable_ocr)
            if result:
                result["format"] = "ppt (partial)"
                return result
        except Exception:
            pass
        return {"error": f"ppt解析失败: {str(e)}", "text": ""}


def _parse_docx(file_path: str, enable_ocr: bool = True) -> dict:
    try:
        from docx import Document
        doc = Document(file_path)

        paragraphs = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                paragraphs.append(text)

        tables_content = []
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data.append(row_data)
            tables_content.append(table_data)

        has_images = False
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                has_images = True
                break

        return {
            "file_path": file_path,
            "format": "docx",
            "total_paragraphs": len(paragraphs),
            "total_tables": len(tables_content),
            "has_images": has_images,
            "paragraphs": paragraphs,
            "tables": tables_content,
            "full_text": "\n".join(paragraphs),
            "image_note": "课件中包含图片，图片内容无法自动提取。如果图片中包含重要文字信息，请用户手动提供图片中的文字内容。" if has_images else None
        }
    except Exception as e:
        return {"error": str(e), "text": ""}


def _parse_doc(file_path: str, enable_ocr: bool = True) -> dict:
    try:
        from src.template_parser.parser import doc_to_docx
        docx_path = doc_to_docx(file_path)
        if docx_path:
            result = _parse_docx(docx_path, enable_ocr)
            result["format"] = "doc (converted)"
            return result
        else:
            return {"error": "doc转docx失败", "text": ""}
    except Exception as e:
        return {"error": str(e), "text": ""}
